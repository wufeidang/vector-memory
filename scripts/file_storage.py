"""
Vector-Memory 记忆系统 - 文件存储模块（混合存储）
负责任意文件的物理存储、文本提取、元数据索引到 ChromaDB
支持：txt/md/csv/pdf/docx/xlsx/png/jpg/gif 等常见格式
"""
import os
import sys
import time
import json
import csv
import hashlib
import mimetypes
import io
from pathlib import Path
from typing import Optional, Dict, Any, List

from core import _get_collection, _get_model, _get_chroma_client, _current_collection_name
from core import MEMORY_MD, CHUNK_SIZE, _sync_to_memory_md, _summarize_if_needed

# ============================================================
# 路径配置
# ============================================================
FILES_DIR = Path(os.path.expanduser("~/.hermes/files"))
FILE_COLLECTION_NAME = "files"

# 支持的文本提取格式
TEXT_FORMATS = {'.txt', '.md', '.csv', '.log', '.json', '.xml', '.html', '.htm', '.sql', '.ini', '.cfg', '.yml', '.yaml', '.toml'}
# 需要额外库的二进制格式
BINARY_TEXT_FORMATS = {'.pdf', '.docx', '.xlsx', '.xls', '.pptx', '.doc', '.xls'}
# 图片格式（存储但用描述文本索引）
IMAGE_FORMATS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.tiff', '.ico'}


# ============================================================
# 工具函数
# ============================================================

def _compute_file_hash(file_path):
    """计算文件内容的 MD5 哈希，用于去重"""
    md5 = hashlib.md5()
    with open(file_path, 'rb') as f:
        for chunk in iter(lambda: f.read(8192), b''):
            md5.update(chunk)
    return md5.hexdigest()


def _get_file_type(file_path):
    """获取文件类型分类"""
    ext = Path(file_path).suffix.lower()
    if ext in TEXT_FORMATS:
        return "text"
    elif ext in BINARY_TEXT_FORMATS:
        return "document"
    elif ext in IMAGE_FORMATS:
        return "image"
    else:
        mime, _ = mimetypes.guess_type(file_path)
        if mime and mime.startswith('text/'):
            return "text"
        return "binary"


def _ensure_file_collection():
    """确保文件元数据集合存在"""
    client = _get_chroma_client()
    existing = [c.name for c in client.list_collections()]
    if FILE_COLLECTION_NAME not in existing:
        client.create_collection(name=FILE_COLLECTION_NAME)
    return _get_collection(FILE_COLLECTION_NAME)


# ============================================================
# 文本提取
# ============================================================

def _extract_text(file_path):
    """
    从文件中提取文本内容。
    根据文件类型使用不同的提取策略。
    
    Returns:
        tuple: (提取的文本, 提取方法, 额外元数据)
    """
    ext = Path(file_path).suffix.lower()
    meta = {}

    # === 纯文本格式 ===
    if ext in TEXT_FORMATS:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            # 限制文本长度（太大时分块处理）
            if len(text) > 50000:
                text = text[:50000] + "\n...[内容截断]"
            return text, "direct_read", meta
        except UnicodeDecodeError:
            # 尝试其他编码
            for enc in ['gbk', 'latin-1', 'cp1252']:
                try:
                    with open(file_path, 'r', encoding=enc) as f:
                        text = f.read()
                    meta["encoding"] = enc
                    if len(text) > 50000:
                        text = text[:50000] + "\n...[内容截断]"
                    return text, "direct_read", meta
                except (UnicodeDecodeError, LookupError):
                    continue
            return "", "failed_encoding", meta

    # === PDF ===
    elif ext == '.pdf':
        text = _extract_pdf_text(file_path)
        meta["page_count"] = _count_pdf_pages(file_path)
        if text:
            return text, "pdf_extract", meta
        return "", "pdf_failed", meta

    # === Word ===
    elif ext in ('.docx', '.doc'):
        text = _extract_docx_text(file_path)
        if text:
            return text, "docx_extract", meta
        return "", "docx_failed", meta

    # === Excel ===
    elif ext in ('.xlsx', '.xls'):
        text = _extract_xlsx_text(file_path)
        if text:
            return text, "xlsx_extract", meta
        return "", "xlsx_failed", meta

    # === PPTX ===
    elif ext == '.pptx':
        text = _extract_pptx_text(file_path)
        if text:
            return text, "pptx_extract", meta
        return "", "pptx_failed", meta

    # === 图片 ===
    elif ext in IMAGE_FORMATS:
        img_meta = _extract_image_metadata(file_path)
        meta.update(img_meta)
        # 图片没有文本内容，用文件名和用户描述作为索引
        return "", "image_no_text", meta

    # === 其他二进制文件 ===
    else:
        return "", "unsupported_format", meta


def _extract_pdf_text(file_path):
    """从 PDF 提取文本"""
    text = ""
    try:
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text() + "\n"
            doc.close()
        except ImportError:
            # 尝试 pypdf
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n"
    except Exception as e:
        print("  ⚠️ PDF提取失败: %s" % str(e)[:80], file=sys.stderr)
    return text.strip()


def _count_pdf_pages(file_path):
    """统计 PDF 页数"""
    try:
        try:
            import fitz
            doc = fitz.open(file_path)
            return len(doc)
        except ImportError:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            return len(reader.pages)
    except Exception:
        return 0


def _extract_docx_text(file_path):
    """从 Word 文档提取文本"""
    text = ""
    try:
        from docx import Document
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + " | "
                text += "\n"
    except ImportError:
        # 尝试用 zipfile 解压读取
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            with zipfile.ZipFile(file_path, 'r') as z:
                with z.open('word/document.xml') as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                    for elem in root.iter():
                        if elem.tag.endswith('}t'):
                            text += elem.text or ''
        except Exception:
            pass
    return text.strip()


def _extract_xlsx_text(file_path):
    """从 Excel 文档提取文本"""
    text = ""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text += f"=== 工作表: {sheet_name} ===\n"
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text += row_text + "\n"
            text += "\n"
        wb.close()
    except ImportError:
        try:
            import xlrd
            wb = xlrd.open_workbook(file_path)
            for sheet in wb.sheets():
                text += f"=== 工作表: {sheet.name} ===\n"
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    row_text = "\t".join(str(cell) for cell in row if cell)
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
        except Exception:
            pass
    return text.strip()


def _extract_pptx_text(file_path):
    """从 PPT 提取文本"""
    text = ""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text += f"=== 幻灯片 {i+1} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            text += "\n"
    except ImportError:
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            with zipfile.ZipFile(file_path, 'r') as z:
                for name in z.namelist():
                    if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                        with z.open(name) as f:
                            tree = ET.parse(f)
                            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                            for elem in tree.iter():
                                if elem.tag.endswith('}t'):
                                    text += (elem.text or '') + "\n"
        except Exception:
            pass
    return text.strip()


def _extract_image_metadata(file_path):
    """提取图片元数据（EXIF等）"""
    meta = {}
    try:
        from PIL import Image
        from PIL.ExifTags import TAGS
        with Image.open(file_path) as img:
            meta["image_size"] = f"{img.width}x{img.height}"
            meta["image_format"] = img.format
            meta["image_mode"] = img.mode
            # 尝试提取 EXIF
            try:
                exif_data = img._getexif()
                if exif_data:
                    for tag_id, value in exif_data.items():
                        tag = TAGS.get(tag_id, tag_id)
                        meta[f"exif_{tag}"] = str(value)[:200]
            except Exception:
                pass
    except ImportError:
        # Pillow 不可用，只记录基本信息
        meta["note"] = "安装 Pillow 可提取更多图片信息"
    except Exception as e:
        meta["error"] = str(e)[:100]
    return meta


# ============================================================
# 文件存储核心 API
# ============================================================

def store_file(args):
    """
    存储文件到混合存储系统。
    
    Args:
        file_path: 本地文件路径
        collection: 目标集合（None = 当前集合）
        description: 用户提供的文件描述（可选，用于图片等无文本内容）
        category: 分类标签
        tags: 标签列表
    
    Returns:
        dict: 操作结果
    """
    file_path = args.get("file_path", "")
    description = args.get("description", "")
    category = args.get("category", "general")
    tags = args.get("tags")
    if tags is None or (isinstance(tags, list) and len(tags) == 0):
        tags = ["uncategorized"]
    collection_name = args.get("collection") or _current_collection_name

    if not file_path or not os.path.exists(file_path):
        return {"success": False, "message": "文件不存在: " + file_path}

    # 1. 计算文件哈希（去重）
    file_hash = _compute_file_hash(file_path)

    # 2. 检查是否已存在
    file_collection = _ensure_file_collection()
    existing = file_collection.get(include=["metadatas"])
    for meta in (existing.get("metadatas") or []):
        if meta.get("file_hash") == file_hash:
            return {
                "success": False,
                "message": "文件已存在（相同内容）",
                "existing_file_id": meta.get("file_id"),
                "file_id": meta.get("file_id")
            }

    # 3. 复制文件到存储目录
    file_id = str(int(time.time() * 1000)) + "_" + hashlib.md5(file_path.encode()).hexdigest()[:8]
    ext = Path(file_path).suffix.lower()
    target_dir = FILES_DIR / collection_name / file_id
    target_dir.mkdir(parents=True, exist_ok=True)
    target_file = target_dir / ("original" + ext)

    import shutil
    shutil.copy2(file_path, target_file)

    # 4. 提取文本
    text, extract_method, extra_meta = _extract_text(file_path)

    # 5. 构建元数据
    file_stat = os.stat(file_path)
    file_size = file_stat.st_size
    filename = Path(file_path).name

    metadata = {
        "file_id": file_id,
        "filename": filename,
        "original_name": filename,
        "file_path": str(target_file),
        "file_type": _get_file_type(file_path),
        "file_ext": ext,
        "file_size": file_size,
        "file_size_human": _human_size(file_size),
        "file_hash": file_hash,
        "extract_method": extract_method,
        "category": category,
        "tags": tags,
        "collection": collection_name,
        "upload_time": time.strftime("%Y-%m-%d %H:%M:%S"),
        "upload_timestamp": time.time(),
        "last_accessed": time.time(),
        "text_length": len(text),
    }
    # 合并额外元数据（图片EXIF等）
    metadata.update(extra_meta)

    # 6. 构建索引文本
    # 对于有文本内容的文件：文件描述 + 提取的文本开头
    # 对于图片/无文本文件：文件描述 + 文件名 + 元数据信息
    if text:
        index_text = (description + " " + text) if description else text
        # 截取用于 embedding 的部分（太长会影响质量）
        index_text_for_embed = index_text[:2000]
    else:
        # 构建描述性文本
        desc_parts = [description] if description else []
        desc_parts.append(filename)
        desc_parts.append("文件类型: " + ext[1:].upper())
        desc_parts.append("文件大小: " + metadata["file_size_human"])
        if extra_meta.get("image_size"):
            desc_parts.append("图片尺寸: " + extra_meta["image_size"])
        index_text = " | ".join(desc_parts)
        index_text_for_embed = index_text

    metadata["description"] = description
    metadata["summary"] = _generate_file_summary(filename, ext, file_size, extract_method, extra_meta)

    # 7. 生成嵌入并存储到 ChromaDB
    model = _get_model()
    if model is None:
        # 清理已复制的文件
        if target_file.exists():
            target_file.unlink()
            target_dir.rmdir()
        return {"success": False, "message": "嵌入模型加载失败"}

    try:
        embedding = model.encode(index_text_for_embed).tolist()
        file_collection.add(
            ids=[file_id],
            embeddings=[embedding],
            documents=[index_text],
            metadatas=[metadata]
        )
    except Exception as e:
        # 嵌入失败，清理文件
        if target_file.exists():
            target_file.unlink()
            target_dir.rmdir()
        return {"success": False, "message": "嵌入生成失败: " + str(e)[:100]}

    # 8. 同步到 MEMORY.md
    _sync_to_memory_md(
        f"[文件] {filename} — {description or '已上传'}",
        {"category": category, "file_id": file_id, "source": "file_upload"}
    )

    return {
        "success": True,
        "file_id": file_id,
        "filename": filename,
        "file_path": str(target_file),
        "file_type": metadata["file_type"],
        "file_size_human": metadata["file_size_human"],
        "text_extracted": len(text) > 0,
        "text_length": len(text),
        "message": "文件已存储并索引"
    }


def list_files(args=None):
    """
    列出所有已索引的文件。
    
    Args:
        limit: 最大数量
        offset: 偏移量
        collection: 过滤集合
        file_type: 过滤文件类型
    
    Returns:
        dict: 文件列表
    """
    try:
        file_collection = _ensure_file_collection()
        data = file_collection.get(include=["metadatas"])
        files = []
        metas = data.get("metadatas") or []

        limit = (args or {}).get("limit", 50)
        offset = (args or {}).get("offset", 0)
        filter_collection = (args or {}).get("collection")
        filter_type = (args or {}).get("file_type")

        for i, meta in enumerate(metas):
            if filter_collection and meta.get("collection") != filter_collection:
                continue
            if filter_type and meta.get("file_type") != filter_type:
                continue

            files.append({
                "file_id": meta.get("file_id"),
                "filename": meta.get("filename"),
                "file_path": meta.get("file_path"),
                "file_type": meta.get("file_type"),
                "file_ext": meta.get("file_ext"),
                "file_size": meta.get("file_size"),
                "file_size_human": meta.get("file_size_human"),
                "description": meta.get("description", ""),
                "category": meta.get("category", "general"),
                "tags": meta.get("tags", []),
                "upload_time": meta.get("upload_time"),
                "text_length": meta.get("text_length", 0),
                "extract_method": meta.get("extract_method", ""),
                "collection": meta.get("collection", ""),
            })

        total = len(files)
        files = files[offset:offset + limit]

        return {
            "success": True,
            "files": files,
            "count": len(files),
            "total": total,
            "offset": offset,
            "limit": limit,
        }
    except Exception as e:
        return {"success": False, "message": str(e)[:100]}


def get_file(args):
    """
    获取文件详情和访问路径。
    
    Args:
        file_id: 文件ID
    
    Returns:
        dict: 文件详情
    """
    try:
        file_collection = _ensure_file_collection()
        data = file_collection.get(ids=[args["file_id"]], include=["metadatas", "documents"])
        if not data.get("ids"):
            return {"success": False, "message": "文件未找到"}

        meta = data["metadatas"][0]
        doc_text = (data.get("documents") or [""])[0]

        result = {
            "file_id": meta.get("file_id"),
            "filename": meta.get("filename"),
            "file_path": meta.get("file_path"),
            "file_type": meta.get("file_type"),
            "file_ext": meta.get("file_ext"),
            "file_size": meta.get("file_size"),
            "file_size_human": meta.get("file_size_human"),
            "description": meta.get("description", ""),
            "category": meta.get("category", "general"),
            "tags": meta.get("tags", []),
            "upload_time": meta.get("upload_time"),
            "extract_method": meta.get("extract_method", ""),
            "text_length": meta.get("text_length", 0),
            "collection": meta.get("collection", ""),
            "indexed_text_preview": doc_text[:500],
        }

        # 添加图片特有信息
        for key in ["image_size", "image_format", "image_mode"]:
            if key in meta:
                result[key] = meta[key]

        return result
    except Exception as e:
        return {"success": False, "message": str(e)[:100]}


def delete_file(args):
    """
    删除文件及其索引。
    
    Args:
        file_id: 文件ID
    
    Returns:
        dict: 操作结果
    """
    try:
        file_collection = _ensure_file_collection()
        data = file_collection.get(ids=[args["file_id"]], include=["metadatas"])
        if not data.get("ids"):
            return {"success": False, "message": "文件未找到"}

        meta = data["metadatas"][0]
        file_path = meta.get("file_path", "")

        # 删除物理文件
        if file_path and os.path.exists(file_path):
            os.remove(file_path)
            # 尝试删除空目录
            parent_dir = os.path.dirname(file_path)
            try:
                if not os.listdir(parent_dir):
                    os.rmdir(parent_dir)
            except Exception:
                pass

        # 删除 ChromaDB 索引
        file_collection.delete(ids=[args["file_id"]])

        return {"success": True, "message": "文件已删除", "file_id": args["file_id"]}
    except Exception as e:
        return {"success": False, "message": str(e)[:100]}


def search_files(args):
    """
    在文件索引中搜索文件。
    
    Args:
        text: 搜索文本
        top_k: 返回数量
        file_type: 按文件类型过滤
        collection: 按集合过滤
    
    Returns:
        dict: 搜索结果
    """
    text = args.get("text", "")
    if not text:
        return {"success": False, "message": "搜索关键词不能为空"}

    top_k = args.get("top_k", 10)
    file_type = args.get("file_type")
    collection_filter = args.get("collection")

    file_collection = _ensure_file_collection()
    model = _get_model()

    try:
        vec = model.encode(text).tolist()
        n_results = top_k * 2
        raw = file_collection.query(
            query_embeddings=[vec],
            n_results=n_results,
            include=["distances", "metadatas", "documents"]
        )
    except Exception as e:
        return {"success": False, "message": "搜索失败: " + str(e)[:100]}

    ids = raw.get("ids", [[]])[0]
    distances = raw.get("distances", [[]])[0]
    documents = raw.get("documents", [[]])[0]
    metadatas = raw.get("metadatas", [[]])[0]

    results = []
    for mem_id, dist, doc, meta in zip(ids, distances, documents, metadatas):
        # 过滤
        if file_type and meta.get("file_type") != file_type:
            continue
        if collection_filter and meta.get("collection") != collection_filter:
            continue

        results.append({
            "file_id": meta.get("file_id"),
            "filename": meta.get("filename"),
            "file_path": meta.get("file_path"),
            "file_type": meta.get("file_type"),
            "file_ext": meta.get("file_ext"),
            "file_size_human": meta.get("file_size_human"),
            "description": meta.get("description", ""),
            "score": 1.0 - (dist / 2.0) if dist else 0.0,
            "text_preview": doc[:200],
            "category": meta.get("category", "general"),
            "upload_time": meta.get("upload_time"),
        })

    # 按分数排序
    results.sort(key=lambda x: x["score"], reverse=True)
    results = results[:top_k]

    # 更新访问时间
    for r in results:
        try:
            meta = file_collection.get(ids=[r["file_id"]], include=["metadatas"])["metadatas"][0]
            meta["last_accessed"] = time.time()
            file_collection.update(ids=[r["file_id"]], metadatas=[meta])
        except Exception:
            pass

    return {
        "success": True,
        "results": results,
        "count": len(results),
        "query": text,
    }


def get_file_stats(args=None):
    """获取文件存储统计信息"""
    file_collection = _ensure_file_collection()
    data = file_collection.get(include=["metadatas"])
    metas = data.get("metadatas") or []

    total_size = sum(m.get("file_size", 0) for m in metas)
    type_counts = {}
    collection_counts = {}

    for m in metas:
        ft = m.get("file_type", "unknown")
        type_counts[ft] = type_counts.get(ft, 0) + 1
        coll = m.get("collection", "unknown")
        collection_counts[coll] = collection_counts.get(coll, 0) + 1

    return {
        "success": True,
        "total_files": len(metas),
        "total_size": total_size,
        "total_size_human": _human_size(total_size),
        "by_type": type_counts,
        "by_collection": collection_counts,
    }


def _human_size(size_bytes):
    """将字节数转换为人类可读格式"""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    elif size_bytes < 1024 * 1024 * 1024:
        return f"{size_bytes / (1024 * 1024):.1f} MB"
    else:
        return f"{size_bytes / (1024 * 1024 * 1024):.2f} GB"


def _generate_file_summary(filename, ext, file_size, extract_method, extra_meta):
    """生成文件的简短摘要描述"""
    parts = [f"文件名: {filename}"]
    parts.append(f"类型: {ext[1:].upper() if ext else '未知'}")
    parts.append(f"大小: {_human_size(file_size)}")

    if extra_meta.get("image_size"):
        parts.append(f"尺寸: {extra_meta['image_size']}")
    if extra_meta.get("page_count"):
        parts.append(f"页数: {extra_meta['page_count']}")
    if extract_method == "image_no_text":
        parts.append("内容: 图片文件（无文本内容）")
    elif extract_method in ("failed_encoding", "pdf_failed", "docx_failed", "xlsx_failed"):
        parts.append("内容: 提取失败（文件可能损坏或格式不支持）")
    elif extract_method != "unsupported_format":
        parts.append("内容: 已成功提取文本")

    return " | ".join(parts)